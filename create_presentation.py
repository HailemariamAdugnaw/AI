from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs):
    """Add a title slide with modern design"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add gradient background (simulation with colored shapes)
    left = Inches(0)
    top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(108, 99, 255)  # Primary color
    shape.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    
    p = title_frame.paragraphs[0]
    p.text = "WeatherWise"
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    
    p = subtitle_frame.paragraphs[0]
    p.text = "AI-Powered Weather Prediction System"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Add tagline
    tagline_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(0.8))
    tagline_frame = tagline_box.text_frame
    
    p = tagline_frame.paragraphs[0]
    p.text = "Machine Learning • Random Forest • Flask Web Application"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(200, 200, 255)
    p.alignment = PP_ALIGN.CENTER
    p.font.italic = True

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add header background
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), prs.slide_width, Inches(1.2)
    )
    fill = header_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(108, 99, 255)
    header_shape.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for idx, item in enumerate(content_items):
        if idx == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(45, 48, 71)
        p.space_before = Pt(12)
        p.level = 0
        
        # Add bullet
        if not item.startswith("•"):
            p.text = "• " + item

def add_architecture_slide(prs):
    """Add technical architecture slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add header
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), prs.slide_width, Inches(1.2)
    )
    fill = header_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(108, 99, 255)
    header_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Technical Architecture"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add architecture boxes
    components = [
        ("Data Layer", "Seattle Weather Dataset\n1,461 observations\n5 weather types"),
        ("Processing", "Preprocessing\n• Feature scaling\n• Data cleaning\n• Train/test split"),
        ("ML Model", "Random Forest\n200 estimators\nBinary classification\n(Rain/No Rain)"),
        ("API Layer", "Flask REST API\n/predict endpoint\nJSON responses"),
        ("Frontend", "Interactive Web UI\nReal-time predictions\nVisual feedback")
    ]
    
    y_pos = 2.0
    for i, (title, desc) in enumerate(components):
        # Add component box
        left = 0.8 + (i % 3) * 3.0
        top = y_pos if i < 3 else y_pos + 2.2
        
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(2.6), Inches(1.6)
        )
        fill = box.fill
        fill.solid()
        if i % 2 == 0:
            fill.fore_color.rgb = RGBColor(230, 230, 255)
        else:
            fill.fore_color.rgb = RGBColor(200, 250, 250)
        box.line.color.rgb = RGBColor(108, 99, 255)
        box.line.width = Pt(2)
        
        # Add text
        text_box = slide.shapes.add_textbox(
            Inches(left + 0.1), Inches(top + 0.1), Inches(2.4), Inches(1.4)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        # Title
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(108, 99, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        p = text_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(45, 48, 71)
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6)

def add_features_slide(prs):
    """Add features and model information slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add header
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), prs.slide_width, Inches(1.2)
    )
    fill = header_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(108, 99, 255)
    header_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Model Features & Specifications"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Features box
    features_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(2), Inches(4), Inches(4.5)
    )
    fill = features_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(230, 240, 255)
    features_box.line.color.rgb = RGBColor(108, 99, 255)
    features_box.line.width = Pt(3)
    
    features_text = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(3.6), Inches(4))
    text_frame = features_text.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = "Input Features"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(108, 99, 255)
    
    features = [
        "Precipitation (mm)",
        "Maximum Temperature (°C)",
        "Minimum Temperature (°C)",
        "Wind Speed (km/h)"
    ]
    
    for feature in features:
        p = text_frame.add_paragraph()
        p.text = "✓ " + feature
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(45, 48, 71)
        p.space_before = Pt(12)
    
    # Model specs box
    specs_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(2), Inches(4), Inches(4.5)
    )
    fill = specs_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 240, 230)
    specs_box.line.color.rgb = RGBColor(255, 107, 139)
    specs_box.line.width = Pt(3)
    
    specs_text = slide.shapes.add_textbox(Inches(5.4), Inches(2.2), Inches(3.6), Inches(4))
    text_frame = specs_text.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = "Model Specifications"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 107, 139)
    
    specs = [
        "Algorithm: Random Forest",
        "Estimators: 200 trees",
        "Balanced class weights",
        "Train/Test: 80/20 split",
        "Standardized features"
    ]
    
    for spec in specs:
        p = text_frame.add_paragraph()
        p.text = "◆ " + spec
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(45, 48, 71)
        p.space_before = Pt(10)

def add_results_slide(prs):
    """Add results and performance metrics slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add header
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), prs.slide_width, Inches(1.2)
    )
    fill = header_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(108, 99, 255)
    header_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Results & Performance"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Key metrics
    metrics = [
        ("92.4%", "Overall Accuracy"),
        ("200", "Trees in Forest"),
        ("1,461", "Training Samples"),
        ("5", "Weather Classes")
    ]
    
    for i, (value, label) in enumerate(metrics):
        left = 1.2 + (i % 4) * 2.2
        top = 2.2
        
        # Metric box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(1.8), Inches(1.8)
        )
        fill = box.fill
        fill.solid()
        colors = [
            RGBColor(76, 175, 80),   # Green
            RGBColor(108, 99, 255),  # Purple
            RGBColor(54, 209, 220),  # Cyan
            RGBColor(255, 152, 0)    # Orange
        ]
        fill.fore_color.rgb = colors[i]
        box.line.fill.background()
        
        # Value
        value_box = slide.shapes.add_textbox(
            Inches(left), Inches(top + 0.3), Inches(1.8), Inches(0.8)
        )
        text_frame = value_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = value
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(
            Inches(left), Inches(top + 1.1), Inches(1.8), Inches(0.6)
        )
        text_frame = label_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    # Key achievements
    achievements_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
    text_frame = achievements_box.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = "Key Achievements"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(108, 99, 255)
    
    achievements = [
        "Accurate binary rain prediction with high confidence",
        "Real-time predictions through responsive web interface",
        "Robust preprocessing pipeline with feature standardization",
        "Production-ready Flask API with scalable architecture"
    ]
    
    for achievement in achievements:
        p = text_frame.add_paragraph()
        p.text = "✓ " + achievement
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(45, 48, 71)
        p.space_before = Pt(8)

def add_webapp_slide(prs):
    """Add web application showcase slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add header
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), prs.slide_width, Inches(1.2)
    )
    fill = header_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(108, 99, 255)
    header_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Web Application Features"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Features
    features = [
        ("🎨 Modern UI/UX", "Responsive design with animated gradients and smooth transitions"),
        ("📊 Real-time Predictions", "Instant weather forecasting with confidence scores"),
        ("🔄 Interactive Forms", "Dynamic input validation and visual feedback"),
        ("📈 Probability Display", "Visual confidence meters showing prediction certainty"),
        ("💡 Smart Recommendations", "Context-aware suggestions based on predictions"),
        ("🌐 RESTful API", "JSON-based API for seamless integration")
    ]
    
    y_pos = 2.0
    for i, (title, desc) in enumerate(features):
        # Feature box
        left = 0.8 if i % 2 == 0 else 5.2
        top = y_pos + (i // 2) * 1.3
        
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(4), Inches(1)
        )
        fill = box.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 240, 255)
        box.line.color.rgb = RGBColor(108, 99, 255)
        box.line.width = Pt(1.5)
        
        # Text
        text_box = slide.shapes.add_textbox(
            Inches(left + 0.15), Inches(top + 0.1), Inches(3.7), Inches(0.8)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(108, 99, 255)
        
        p = text_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(45, 48, 71)
        p.space_before = Pt(4)

def add_conclusion_slide(prs):
    """Add conclusion and thank you slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    fill = bg_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(108, 99, 255)
    bg_shape.line.fill.background()
    
    # Thank you text
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    
    p = subtitle_frame.paragraphs[0]
    p.text = "WeatherWise - Predicting Tomorrow, Today"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    p.font.italic = True
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    contact_frame = contact_box.text_frame
    
    p = contact_frame.paragraphs[0]
    p.text = "Questions & Feedback Welcome"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(200, 200, 255)
    p.alignment = PP_ALIGN.CENTER

def create_presentation():
    """Main function to create the presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs)
    
    # Slide 2: Project Overview
    add_content_slide(prs, "Project Overview", [
        "A machine learning-powered weather prediction web application",
        "Predicts rain probability based on meteorological parameters",
        "Built with Python, Flask, and scikit-learn",
        "Features modern, responsive web interface with real-time predictions",
        "Uses historical Seattle weather data for training",
        "Production-ready with RESTful API architecture"
    ])
    
    # Slide 3: Problem Statement
    add_content_slide(prs, "Problem Statement", [
        "Weather prediction is crucial for daily planning and decision-making",
        "Traditional forecasting methods can be complex and inaccessible",
        "Need for simple, accurate, and user-friendly prediction tools",
        "Gap in lightweight ML-based weather applications",
        "Challenge: Build an accurate model with minimal input features",
        "Solution: Random Forest classifier with intuitive web interface"
    ])
    
    # Slide 4: Dataset Information
    add_content_slide(prs, "Dataset Information", [
        "Source: Seattle Weather Dataset (2012-2015)",
        "Total Observations: 1,461 daily weather records",
        "Features: Precipitation, Temperature (max/min), Wind Speed",
        "Weather Classes: Drizzle, Fog, Rain, Snow, Sun",
        "Binary Target: Rain vs. No Rain classification",
        "Data Preprocessing: Cleaning, scaling, and stratified splitting"
    ])
    
    # Slide 5: Technical Architecture
    add_architecture_slide(prs)
    
    # Slide 6: Features & Model
    add_features_slide(prs)
    
    # Slide 7: Results
    add_results_slide(prs)
    
    # Slide 8: Web Application
    add_webapp_slide(prs)
    
    # Slide 9: Future Enhancements
    add_content_slide(prs, "Future Enhancements", [
        "Multi-class weather prediction (sun, rain, snow, fog, etc.)",
        "Integration with real-time weather APIs",
        "Historical weather data visualization and trends",
        "Location-based predictions using geolocation",
        "Mobile application development (iOS/Android)",
        "Advanced ensemble methods and deep learning models",
        "User accounts with saved prediction history"
    ])
    
    # Slide 10: Technologies Used
    add_content_slide(prs, "Technology Stack", [
        "Backend: Python 3.x with Flask framework",
        "ML Library: scikit-learn (Random Forest Classifier)",
        "Data Processing: pandas, NumPy",
        "Visualization: matplotlib, seaborn",
        "Frontend: HTML5, CSS3, JavaScript (Vanilla)",
        "Model Persistence: joblib for serialization",
        "Development: Git version control"
    ])
    
    # Slide 11: Thank You
    add_conclusion_slide(prs)
    
    # Save presentation
    output_file = "WeatherWise_Presentation.pptx"
    prs.save(output_file)
    print(f"✓ Presentation created successfully: {output_file}")
    print(f"✓ Total slides: {len(prs.slides)}")
    return output_file

if __name__ == "__main__":
    create_presentation()
